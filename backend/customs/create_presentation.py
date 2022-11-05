from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from datetime import datetime


# вместо 'img1.png' и 'img2.jpg' долны быть графики по дефолту


def text_analytics(import_custom_volume, export_custom_volume,
                   clean_import, clean_exp_imp_del, main_partners, import_duties, sanctions_import):
    """Для 4-5 слайда"""

    ans_1 = f'Объём импорта за период: {import_custom_volume} тыс. $ ' \
            f'Объём экспорта за период: {export_custom_volume} тыс. $ ' \
            f'\n' \
            f'Чистый импорт (экспорт, если -): {clean_import} ' \
            f' \n' \
            f'Изменение чистого импорта(экспорта, если -): {clean_exp_imp_del}% '
    ans_2 = f'Основные партнеры по импорту: {", ".join([str(x) for x in [*main_partners]])} \n' \
            f'\n' \
            f'Таможенные пошлины на импорт: {import_duties} ' \
            f'Санкции: {", ".join([str(x) for x in [*sanctions_import]])}  \n' \
            f'Потенциальный объем ниши: {clean_import} тыс. $   \n' \
            f'Рост ниши за год: {clean_exp_imp_del} %'

    return [ans_1, ans_2]


subtitle_text_4_foo_default = text_analytics(0, 0, 0, 0,[], 0, [])


def create_presentation(tnved_code='Вы не выбрали код', product="Вы не выбрали товар", img_path_1='default_image.png',
                        img_path_2='defoult_image.jpg',
                        subtitle_2='Чистый импорт в товара (помесячно), тыс.дол.', subtitle_3='Доля подсанкционного импорта, %', title_text_4='Информацоннная справка',
                        subtitle_text_4_foo=subtitle_text_4_foo_default, author='Сотрудник департамента'):
    """tnved_code: product tnved code, product: product name,
       img_path_1: rectangular diagram, img_path_2: square diagram,
       subtitle_2: caption to a rectangular diagram,
       subtitle_3: caption to the square diagram,
       title_text_4: name of text analytics,
       subtitle_text_4_foo: function for the text analytics template,
       author: author name.

    """

    prs = Presentation()

    # first slide
    first_slide_layout = prs.slide_layouts[0]
    first_slide = prs.slides.add_slide(first_slide_layout)

    title = first_slide.shapes.title
    subtitle = first_slide.placeholders[1]

    title.text = "Аналитика по коду ТН ВЭД:"

    p = title.text_frame.add_paragraph()
    p.text = f"{tnved_code}"
    p.font.size = Pt(36)
    p.font.bold = True

    p = subtitle.text_frame.add_paragraph()
    p.text = f"{product}"
    p.font.size = Pt(18)

    # second slide
    second_slide_layout = prs.slide_layouts[6]
    second_slide = prs.slides.add_slide(second_slide_layout)

    # title
    left_title = Inches(0.492)
    top_title = Inches(0.25)
    height_title = Inches(1.25)
    width_title = Inches(9.031496)

    txBox = second_slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
    tf = txBox.text_frame
    tf.text = " "

    p = tf.add_paragraph()
    p.text = "Объёмы внешней торговли"
    p.font.size = Pt(44)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # picture
    left_1 = Inches(2.28)
    top = Inches(1.94)
    height = Inches(3.41)

    width = height * 1.6

    pic_1 = second_slide.shapes.add_picture(img_path_1, left_1, top, height=height, width=width)

    # subtitle

    left_subtitle = Inches(1.2677)
    top_subtitle = height + top + Inches(0.25)
    height_subtitle = Inches(0.748)
    width_subtitle = Inches(7.492)

    txBox_2 = second_slide.shapes.add_textbox(left_subtitle, top_subtitle, width_subtitle, height_subtitle)
    tf_2 = txBox_2.text_frame

    tf_2.text = " "
    p_2 = tf_2.add_paragraph()
    p_2.text = f"{subtitle_2}"
    p_2.font.size = Pt(20)
    p_2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # third slide
    third_slide_layout = prs.slide_layouts[6]
    third_slide = prs.slides.add_slide(third_slide_layout)

    # title
    left_title = Inches(0.492)
    top_title = Inches(0.25)
    height_title = Inches(1.25)
    width_title = Inches(9.031496)

    txBox = third_slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
    tf = txBox.text_frame
    tf.text = " "

    p = tf.add_paragraph()
    p.text = "Торговые партнеры"
    p.font.size = Pt(44)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # picture
    left_1 = Inches(3.3228)
    top = Inches(1.94)
    height = Inches(3.41)

    width = height

    pic_2 = third_slide.shapes.add_picture(img_path_2, left_1, top, height=height, width=width)

    # subtitle

    left_subtitle = Inches(1.2677)
    top_subtitle = height + top + Inches(0.25)
    height_subtitle = Inches(0.748)
    width_subtitle = Inches(7.492)

    txBox_2 = third_slide.shapes.add_textbox(left_subtitle, top_subtitle, width_subtitle, height_subtitle)
    tf_2 = txBox_2.text_frame

    tf_2.text = " "
    p_2 = tf_2.add_paragraph()
    p_2.text = f"{subtitle_3}"
    p_2.font.size = Pt(20)
    p_2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # fourth slide
    fourth_slide_layout = prs.slide_layouts[1]
    fourth_slide = prs.slides.add_slide(fourth_slide_layout)

    title = fourth_slide.shapes.title
    subtitle = fourth_slide.placeholders[1]

    title.text = f"{title_text_4}"
    subtitle.text = subtitle_text_4_foo[0]

    # fourth slide 4.5
    fourth_slide_layout = prs.slide_layouts[1]
    fourth_slide = prs.slides.add_slide(fourth_slide_layout)

    title = fourth_slide.shapes.title
    subtitle = fourth_slide.placeholders[1]

    title.text = f"{title_text_4}"
    subtitle.text = subtitle_text_4_foo[1]

    # fifth slide
    fifth_slide_layout = prs.slide_layouts[0]
    fifth_slide = prs.slides.add_slide(fifth_slide_layout)

    title = fifth_slide.shapes.title
    subtitle = fifth_slide.placeholders[1]

    title.text = f"Автор презентации: \n {author}"
    current_datetime = datetime.now()
    subtitle.text = f"{current_datetime.day}.{current_datetime.month}.{current_datetime.year} " \
                    f"\n Московское Таможенное Управление"

    # save presentation
    prs.save('test.pptx')

# if __name__ == "__main__":
#     create_presentation()
#     print(text_analytics('штук', 1000, 1200, 123, 45654, 'dfggdgf', 456456, 56456456465, ["ыавп", "рвыаолр"], ["ыавп", "рвыаолр"],
#                    ["ыавп", "рвыаолр"], ["ыавп", "рвыаолр"], 'rjhdgidshgisudfghlsdfkjghdfskjghdfskl', "becouse sdjkgfhdskjlfhsdkjl"))
